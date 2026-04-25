"""
PilotGuard — Presentation Generator
Generates a polished .pptx file for a 5-minute video presentation.
Run: pip install python-pptx Pillow  &&  python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD    = RGBColor(0x16, 0x21, 0x3E)   # card navy
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)   # bright cyan
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)   # teal
GREEN      = RGBColor(0x4A, 0xDE, 0x80)   # green
AMBER      = RGBColor(0xFB, 0xBF, 0x24)   # amber
RED        = RGBColor(0xF8, 0x71, 0x71)   # red/coral
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ════════════════════════════════════════════════════════════════
#  HELPERS
# ════════════════════════════════════════════════════════════════

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txbox = slide.shapes.add_textbox(left, top, w, h)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txbox

def add_para(tf, text, size=16, color=LIGHT_GRAY, bold=False, space_before=Pt(6), bullet=False, font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p

def slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num}/{total}", size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def accent_line(slide, left, top, width):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return line

def section_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             title, size=32, color=WHITE, bold=True)
    accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
                 subtitle, size=16, color=MID_GRAY)

TOTAL_SLIDES = 12

# ════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl, BG_DARK)

# Large accent circle (decorative)
circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.5), Inches(6), Inches(6))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x1E, 0x30, 0x50)
circ.line.fill.background()

circ2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
circ2.fill.solid()
circ2.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x48)
circ2.line.fill.background()

# "PG" badge
badge = add_rect(sl, Inches(1.0), Inches(1.8), Inches(1.2), Inches(1.2), ACCENT, radius=0.15)
badge_tf = badge.text_frame
badge_tf.paragraphs[0].text = "PG"
badge_tf.paragraphs[0].font.size = Pt(36)
badge_tf.paragraphs[0].font.bold = True
badge_tf.paragraphs[0].font.color.rgb = BG_DARK
badge_tf.paragraphs[0].font.name = "Segoe UI"
badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
badge_tf.word_wrap = False

add_text(sl, Inches(1.0), Inches(3.2), Inches(10), Inches(1.2),
         "PilotGuard", size=52, color=WHITE, bold=True)
add_text(sl, Inches(1.0), Inches(4.2), Inches(10), Inches(0.8),
         "Real-Time Pilot Cognitive State Monitoring", size=28, color=ACCENT)
add_text(sl, Inches(1.0), Inches(5.0), Inches(10), Inches(0.6),
         "Computer Vision  |  Deep Learning  |  Signal Processing", size=18, color=LIGHT_GRAY)

accent_line(sl, Inches(1.0), Inches(5.7), Inches(4))

add_text(sl, Inches(1.0), Inches(5.9), Inches(6), Inches(0.4),
         "Om Patel  —  Northeastern University  —  April 2026", size=15, color=MID_GRAY)
add_text(sl, Inches(1.0), Inches(6.3), Inches(8), Inches(0.4),
         "github.com/omunite215/Project_Pilotguard", size=13, color=ACCENT2)
slide_number(sl, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 2 — MOTIVATION / PROBLEM
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Why PilotGuard?", "The problem that drives this project")

# Stat cards
stats = [
    ("70-80%", "of aviation accidents\ninvolve human factors", RED),
    ("#1 Factor", "Pilot fatigue is the\nleading contributor", AMBER),
    ("Unreliable", "Self-assessment of\nfatigue is subjective", PURPLE),
]
for i, (num, desc, col) in enumerate(stats):
    x = Inches(0.8 + i * 3.9)
    card = add_rect(sl, x, Inches(2.0), Inches(3.5), Inches(2.0), BG_CARD, border_color=col, radius=0.08)
    add_text(sl, x + Inches(0.3), Inches(2.2), Inches(2.9), Inches(0.8),
             num, size=36, color=col, bold=True)
    add_text(sl, x + Inches(0.3), Inches(3.0), Inches(2.9), Inches(0.8),
             desc, size=15, color=LIGHT_GRAY)

# Solution box
sol = add_rect(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.3), RGBColor(0x0A, 0x2E, 0x1A), border_color=GREEN, radius=0.05)
add_text(sl, Inches(1.2), Inches(4.7), Inches(3), Inches(0.5),
         "The Solution", size=22, color=GREEN, bold=True)
add_text(sl, Inches(1.2), Inches(5.2), Inches(10.8), Inches(1.4),
         "PilotGuard acts as a non-intrusive \"third crew member\" — a webcam-based system that continuously "
         "monitors the pilot's face for signs of drowsiness, fatigue, and microsleep events, delivering "
         "graded alerts before cognitive impairment reaches dangerous levels. No wearables. No EEG. Just a camera.",
         size=16, color=LIGHT_GRAY)
slide_number(sl, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 3 — TOOLS & TECH STACK
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Tools & Technology Stack", "What powers PilotGuard")

cols = [
    ("Computer Vision", ACCENT, [
        "MediaPipe FaceLandmarker (478 pts)",
        "OpenCV (frame decode/preprocess)",
        "Kalman Filter (signal smoothing)",
        "PERCLOS + Adaptive EAR",
    ]),
    ("Machine Learning", GREEN, [
        "PyTorch 2.x + CUDA",
        "DINOv2 ViT-S/14 (transfer learning)",
        "XGBoost (geometric classifier)",
        "HMM + Bayesian confidence",
    ]),
    ("Backend", AMBER, [
        "FastAPI + Uvicorn",
        "WebSocket (real-time frames)",
        "SQLite (session persistence)",
        "Pydantic (typed API schemas)",
    ]),
    ("Frontend", PURPLE, [
        "React 19 + TypeScript",
        "Zustand (state management)",
        "TanStack Router + Query",
        "Tailwind CSS v4 + Motion",
    ]),
]

for i, (title, col, items) in enumerate(cols):
    x = Inches(0.5 + i * 3.1)
    card = add_rect(sl, x, Inches(1.9), Inches(2.9), Inches(4.8), BG_CARD, border_color=col, radius=0.06)
    # Color bar at top
    bar = add_rect(sl, x, Inches(1.9), Inches(2.9), Pt(5), col, radius=0.0)
    add_text(sl, x + Inches(0.2), Inches(2.15), Inches(2.5), Inches(0.5),
             title, size=18, color=col, bold=True)
    for j, item in enumerate(items):
        add_text(sl, x + Inches(0.25), Inches(2.8 + j * 0.55), Inches(2.5), Inches(0.5),
                 f"  {item}", size=13, color=LIGHT_GRAY)
        # Small dot
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(2.93 + j * 0.55), Pt(7), Pt(7))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.fill.background()

# HPC callout
hpc = add_rect(sl, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45), RGBColor(0x1A, 0x1A, 0x3E), border_color=ACCENT2, radius=0.05)
add_text(sl, Inches(0.8), Inches(6.87), Inches(11.5), Inches(0.4),
         "Training Infrastructure:  NEU Discovery HPC  |  NVIDIA H200 GPU  |  SLURM scheduler  |  Mixed-precision bf16",
         size=13, color=ACCENT2)
slide_number(sl, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 4 — SYSTEM ARCHITECTURE (image)
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "System Architecture", "Four-layer design from camera to dashboard")

img_path = os.path.join(os.path.dirname(__file__), "system_architecture.png")
if os.path.exists(img_path):
    # Center the image
    sl.shapes.add_picture(img_path, Inches(0.4), Inches(1.6), Inches(12.5), Inches(5.6))
else:
    add_text(sl, Inches(2), Inches(3), Inches(9), Inches(2),
             "[Place system_architecture.png here]", size=28, color=MID_GRAY, align=PP_ALIGN.CENTER)

slide_number(sl, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 5 — CV PIPELINE (deep dive, focus slide)
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Computer Vision Pipeline", "The core — 7 stages, under 30ms per frame")

steps = [
    ("1", "Preprocess", "Resize 640x480\nBGR to RGB", ACCENT),
    ("2", "Face Detect", "MediaPipe\n478 landmarks", ACCENT),
    ("3", "EAR Compute", "6 eye landmarks\nper-eye ratio", GREEN),
    ("4", "Kalman Filter", "Smooth signal\nnoise removal", GREEN),
    ("5", "Blink Detect", "Adaptive thresh\n0.75x baseline", AMBER),
    ("6", "PERCLOS", "60s rolling\nwindow", AMBER),
    ("7", "Classify State", "Alert / Drowsy\n/ Microsleep", RED),
]

for i, (num, title, desc, col) in enumerate(steps):
    x = Inches(0.3 + i * 1.82)
    # Step box
    card = add_rect(sl, x, Inches(2.0), Inches(1.65), Inches(2.5), BG_CARD, border_color=col, radius=0.08)
    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(2.15), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    circ_tf = circ.text_frame
    circ_tf.paragraphs[0].text = num
    circ_tf.paragraphs[0].font.size = Pt(18)
    circ_tf.paragraphs[0].font.bold = True
    circ_tf.paragraphs[0].font.color.rgb = BG_DARK
    circ_tf.paragraphs[0].font.name = "Segoe UI"
    circ_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(sl, x + Inches(0.1), Inches(2.8), Inches(1.45), Inches(0.4),
             title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), Inches(3.2), Inches(1.45), Inches(0.8),
             desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < 6:
        arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.68), Inches(3.05), Inches(0.18), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# EAR formula box
ear_box = add_rect(sl, Inches(0.5), Inches(4.8), Inches(5.5), Inches(2.2), BG_CARD, border_color=ACCENT, radius=0.06)
add_text(sl, Inches(0.8), Inches(4.95), Inches(5), Inches(0.4),
         "Eye Aspect Ratio (EAR)", size=18, color=ACCENT, bold=True)
add_text(sl, Inches(0.8), Inches(5.4), Inches(5), Inches(0.6),
         "EAR = (||p2-p6|| + ||p3-p5||) / (2 * ||p1-p4||)", size=16, color=WHITE, font_name="Consolas")
add_text(sl, Inches(0.8), Inches(5.95), Inches(5), Inches(0.9),
         "Computed from 6 landmarks per eye. Drops to ~0 during blinks.\n"
         "Adaptive threshold: 0.75 x calibrated baseline mean (30s calibration).\n"
         "Kalman filter: process noise 0.01, measurement noise 0.1.",
         size=12, color=LIGHT_GRAY)

# State classification box
state_box = add_rect(sl, Inches(6.4), Inches(4.8), Inches(6.4), Inches(2.2), BG_CARD, border_color=RED, radius=0.06)
add_text(sl, Inches(6.7), Inches(4.95), Inches(5.5), Inches(0.4),
         "State Classification Logic", size=18, color=RED, bold=True)

states_info = [
    ("ALERT", "PERCLOS < 40% AND EAR above threshold", GREEN),
    ("DROWSY", "PERCLOS > 60% OR EAR < 80% of baseline", AMBER),
    ("MICROSLEEP", "Eyes closed > 1.5 seconds continuously", RED),
]
for j, (state, cond, col) in enumerate(states_info):
    y = Inches(5.45 + j * 0.5)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.7), y + Inches(0.08), Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    add_text(sl, Inches(7.05), y, Inches(1.2), Inches(0.35),
             state, size=13, color=col, bold=True)
    add_text(sl, Inches(8.2), y, Inches(4.2), Inches(0.35),
             cond, size=12, color=LIGHT_GRAY)

slide_number(sl, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 6 — DINOv2 + ML MODELS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Deep Learning Models", "DINOv2 transfer learning + geometric fusion")

# DINOv2 card
dino_card = add_rect(sl, Inches(0.5), Inches(1.9), Inches(6), Inches(3.0), BG_CARD, border_color=ACCENT, radius=0.06)
add_text(sl, Inches(0.8), Inches(2.05), Inches(5.5), Inches(0.4),
         "DINOv2 ViT-S/14  (Self-Supervised)", size=20, color=ACCENT, bold=True)
txb = add_text(sl, Inches(0.8), Inches(2.55), Inches(5.5), Inches(2.2),
               "", size=14, color=LIGHT_GRAY)
tf = txb.text_frame
tf.paragraphs[0].text = "384-dimensional feature vectors from frozen backbone"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].font.name = "Segoe UI"
add_para(tf, "Trained on 142M images — captures texture & expression", size=14, color=LIGHT_GRAY)
add_para(tf, "Runs every 3rd frame (cached) to save compute", size=14, color=LIGHT_GRAY)
add_para(tf, "Input: 224x224 center-cropped face region", size=14, color=LIGHT_GRAY)
add_para(tf, "No fine-tuning of backbone — only head is trained", size=14, color=AMBER)

# Three head architectures
heads = [
    ("LinearProbe", "LayerNorm + Linear\nNTHU drowsiness", ACCENT),
    ("MLPProbe", "3-layer MLP + GELU\nEmotion & AU tasks", GREEN),
    ("FusionHead", "384 deep + 7 geometric\nUTA drowsiness", AMBER),
]
for i, (name, desc, col) in enumerate(heads):
    x = Inches(0.5 + i * 4.2)
    card = add_rect(sl, x, Inches(5.2), Inches(3.8), Inches(1.6), BG_CARD, border_color=col, radius=0.06)
    add_text(sl, x + Inches(0.2), Inches(5.3), Inches(3.4), Inches(0.35),
             name, size=16, color=col, bold=True)
    add_text(sl, x + Inches(0.2), Inches(5.7), Inches(3.4), Inches(0.9),
             desc, size=13, color=LIGHT_GRAY)

# XGBoost card
xgb_card = add_rect(sl, Inches(6.8), Inches(1.9), Inches(5.8), Inches(3.0), BG_CARD, border_color=GREEN, radius=0.06)
add_text(sl, Inches(7.1), Inches(2.05), Inches(5.2), Inches(0.4),
         "XGBoost  (Geometric Baseline)", size=20, color=GREEN, bold=True)
txb2 = add_text(sl, Inches(7.1), Inches(2.55), Inches(5.2), Inches(2.2),
                "", size=14, color=LIGHT_GRAY)
tf2 = txb2.text_frame
tf2.paragraphs[0].text = "7 hand-crafted features per frame:"
tf2.paragraphs[0].font.size = Pt(14)
tf2.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf2.paragraphs[0].font.name = "Segoe UI"
add_para(tf2, "  EAR (L/R/Avg), MAR, Blink Rate, PERCLOS, EAR dev", size=13, color=WHITE)
add_para(tf2, "1000 estimators, max depth 6, lr 0.05", size=13, color=LIGHT_GRAY)
add_para(tf2, "Fast CPU inference — no GPU needed", size=13, color=LIGHT_GRAY)
add_para(tf2, "Serves as fallback when DINOv2 too slow", size=13, color=AMBER)

slide_number(sl, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 7 — DATA PIPELINE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Data Pipeline", "Four datasets, seven-step cleaning, subject-stratified splits")

datasets = [
    ("NTHU-DDD", "68,564 images", "Binary drowsiness\n(alert vs drowsy)", ACCENT),
    ("UTA-RLDD", "125,490 frames", "Naturalistic drowsiness\n(60 subjects, 30h video)", GREEN),
    ("AffectNet", "19,328 images", "5 emotions (neutral,\nstress, surprise...)", AMBER),
    ("DISFA", "36,080 images", "Action Unit labels\n(12 AUs, multi-label)", PURPLE),
]
for i, (name, count, desc, col) in enumerate(datasets):
    x = Inches(0.5 + i * 3.15)
    card = add_rect(sl, x, Inches(1.9), Inches(2.9), Inches(1.8), BG_CARD, border_color=col, radius=0.06)
    add_text(sl, x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.35),
             name, size=16, color=col, bold=True)
    add_text(sl, x + Inches(0.2), Inches(2.35), Inches(2.5), Inches(0.3),
             count, size=13, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.8),
             desc, size=12, color=LIGHT_GRAY)

# Cleaning pipeline
clean_box = add_rect(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0), BG_CARD, border_color=ACCENT, radius=0.06)
add_text(sl, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
         "7-Step Cleaning Pipeline", size=18, color=ACCENT, bold=True)

steps_clean = [
    ("1. Corrupt detection", "Remove broken/unreadable images"),
    ("2. Quality checks", "Resolution + quality score filtering"),
    ("3. CLAHE enhancement", "Contrast normalization for low-light"),
    ("4. Label remapping", "Map to task-specific label space"),
    ("5. Class balancing", "Undersample majority classes"),
    ("6. Subject-stratified split", "70/15/15 — NO subject leakage"),
    ("7. Manifest generation", "CSV with forward-slash paths (cross-platform)"),
]
for j, (step, desc) in enumerate(steps_clean):
    col_idx = j % 2
    x_base = Inches(0.8 + col_idx * 6.2)
    y_base = Inches(4.7 + (j // 2) * 0.55)
    add_text(sl, x_base, y_base, Inches(2.2), Inches(0.4),
             step, size=13, color=WHITE, bold=True)
    add_text(sl, x_base + Inches(2.2), y_base, Inches(3.8), Inches(0.4),
             desc, size=12, color=LIGHT_GRAY)

slide_number(sl, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 8 — ALERT ENGINE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Alert Engine", "Safety-critical tiered alerting with debounce and lock")

# Three alert level cards
levels = [
    ("ADVISORY", "Score >= 30", "15s lock\n15 frames to unlock", AMBER, "Mild drowsiness detected.\nReminder to stay alert."),
    ("CAUTION", "Score >= 55", "30s lock\n10 frames to unlock", RGBColor(0xF9, 0x73, 0x16), "Significant drowsiness.\nConsider a break soon."),
    ("WARNING", "Score >= 75", "60s lock\n5 frames to unlock", RED, "Severe impairment.\nImmediate action needed."),
]
for i, (level, thresh, lock, col, msg) in enumerate(levels):
    x = Inches(0.5 + i * 4.1)
    card = add_rect(sl, x, Inches(2.0), Inches(3.8), Inches(2.8), BG_CARD, border_color=col, radius=0.06)
    bar = add_rect(sl, x, Inches(2.0), Inches(3.8), Pt(5), col, radius=0.0)
    add_text(sl, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.4),
             level, size=22, color=col, bold=True)
    add_text(sl, x + Inches(0.3), Inches(2.65), Inches(3.2), Inches(0.3),
             thresh, size=14, color=WHITE, bold=True)
    add_text(sl, x + Inches(0.3), Inches(3.0), Inches(3.2), Inches(0.6),
             lock, size=13, color=LIGHT_GRAY)
    add_text(sl, x + Inches(0.3), Inches(3.7), Inches(3.2), Inches(0.8),
             msg, size=12, color=MID_GRAY)

# Key features
feat_box = add_rect(sl, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), BG_CARD, border_color=PURPLE, radius=0.06)
add_text(sl, Inches(0.8), Inches(5.3), Inches(5), Inches(0.4),
         "Safety Mechanisms", size=18, color=PURPLE, bold=True)

feats = [
    ("Debounce:", "8 consecutive frames for microsleep, 15 for drowsiness — eliminates false alarms"),
    ("Mandatory Lock:", "Pilot cannot dismiss alert until sustained alertness is demonstrated"),
    ("Escalation:", "If microsleep occurs during an active lock, severity escalates to WARNING"),
    ("Audio:", "Looping failure tone during lock (cannot mute), spoken \"All clear\" on unlock"),
]
for j, (label, desc) in enumerate(feats):
    y = Inches(5.7 + j * 0.3)
    add_text(sl, Inches(1.0), y, Inches(1.6), Inches(0.3),
             label, size=12, color=AMBER, bold=True)
    add_text(sl, Inches(2.6), y, Inches(9.5), Inches(0.3),
             desc, size=12, color=LIGHT_GRAY)

slide_number(sl, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 9 — FRONTEND DASHBOARD
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Frontend Dashboard", "React 19 real-time monitoring interface")

components = [
    ("VideoFeed + Landmarks", "Live camera with 478-point\noverlay, state banner,\nemotion badge, latency", ACCENT, Inches(0.5), Inches(2.0), Inches(3.8), Inches(2.5)),
    ("FatigueGauge (SVG)", "Circular gauge, 4 color\nzones, spring animations,\nglow effects", GREEN, Inches(4.6), Inches(2.0), Inches(3.8), Inches(2.5)),
    ("Session Timeline", "Dual-axis canvas chart\nEAR + Fatigue, 60s window\nthrottled to 1s updates", AMBER, Inches(8.7), Inches(2.0), Inches(4.1), Inches(2.5)),
    ("Alert Banner + Feed", "Full-width alert with\nlock countdown, progress\nbar, dismiss button", RED, Inches(0.5), Inches(4.8), Inches(3.8), Inches(2.0)),
    ("EAR Indicator", "Bar chart with threshold\nmarker, status badge\n(open/closed/calibrating)", PURPLE, Inches(4.6), Inches(4.8), Inches(3.8), Inches(2.0)),
    ("Audio System", "Web Audio API alerts\nSpeech Synthesis voice\n8s cooldown, lock loop", ACCENT2, Inches(8.7), Inches(4.8), Inches(4.1), Inches(2.0)),
]
for name, desc, col, x, y, w, h in components:
    card = add_rect(sl, x, y, w, h, BG_CARD, border_color=col, radius=0.06)
    add_text(sl, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.35),
             name, size=14, color=col, bold=True)
    add_text(sl, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), h - Inches(0.6),
             desc, size=12, color=LIGHT_GRAY)

slide_number(sl, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 10 — CHALLENGES & LESSONS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Challenges & Lessons Learned", "What broke, what worked, what changed")

challenges = [
    ("MediaPipe API Breaking Change",
     "MediaPipe 0.10.9 deprecated mp.solutions.face_mesh entirely. Had to migrate to the new Tasks API "
     "(mp.tasks.vision.FaceLandmarker) which requires a separate .task model file.",
     "Caught during testing. Migration took effort but gave us 478-point landmarks with iris tracking.",
     RED),
    ("Fixed EAR Threshold Fails Across Users",
     "The literature suggests EAR ~ 0.21 as a blink threshold, but users with narrow or wide eyes were "
     "getting constant false positives or missed detections.",
     "Implemented adaptive thresholding: 30s calibration period, threshold = 0.75x user's baseline mean.",
     AMBER),
    ("Subject Leakage Inflates Metrics",
     "Random train/test split allowed the model to memorize individual faces. Validation F1 looked great "
     "but the model could not generalize to new subjects.",
     "Switched to subject-stratified splitting. F1 dropped 5-8% but now reflects real generalization.",
     GREEN),
    ("Cross-Platform Paths on HPC",
     "Windows backslash paths in data manifests failed silently on the Linux HPC cluster. Training jobs "
     "reported 'file not found' with no obvious cause.",
     "All manifest writers now use Path.as_posix(). Rule codified in lessons log.",
     ACCENT),
]

for i, (title, problem, solution, col) in enumerate(challenges):
    y = Inches(1.8 + i * 1.35)
    card = add_rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.2), BG_CARD, border_color=col, radius=0.05)
    add_text(sl, Inches(0.8), y + Inches(0.05), Inches(5), Inches(0.3),
             title, size=14, color=col, bold=True)
    add_text(sl, Inches(0.8), y + Inches(0.35), Inches(6), Inches(0.8),
             problem, size=11, color=LIGHT_GRAY)
    # Solution on the right
    add_text(sl, Inches(7.2), y + Inches(0.05), Inches(1.2), Inches(0.3),
             "Resolution:", size=11, color=GREEN, bold=True)
    add_text(sl, Inches(7.2), y + Inches(0.35), Inches(5.3), Inches(0.8),
             solution, size=11, color=LIGHT_GRAY)

slide_number(sl, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 11 — DEMO (placeholder + instructions)
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)

# Big centered text
add_text(sl, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.0),
         "Live Demonstration", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
accent_line(sl, Inches(5.2), Inches(1.9), Inches(3))

# Demo steps
demo_box = add_rect(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(4.2), BG_CARD, border_color=ACCENT, radius=0.06)

demo_steps = [
    ("1.  Start the backend", "uvicorn backend.src.api.main:app --reload", ACCENT),
    ("2.  Start the frontend", "cd frontend && bun dev", ACCENT),
    ("3.  Open dashboard", "Navigate to http://localhost:5173", GREEN),
    ("4.  Start monitoring session", "Click 'Start Monitoring' — 30s calibration begins", GREEN),
    ("5.  Show landmark overlay", "478 points tracked in real-time on face", AMBER),
    ("6.  Simulate drowsiness", "Close eyes slowly — watch EAR drop, fatigue rise", RED),
    ("7.  Trigger alert", "Sustained eye closure triggers Advisory then Warning", RED),
]

for j, (step, detail, col) in enumerate(demo_steps):
    y = Inches(2.7 + j * 0.53)
    add_text(sl, Inches(1.8), y, Inches(3.2), Inches(0.4),
             step, size=14, color=col, bold=True)
    add_text(sl, Inches(5.0), y, Inches(6.2), Inches(0.4),
             detail, size=13, color=LIGHT_GRAY)

add_text(sl, Inches(1.5), Inches(6.85), Inches(10.3), Inches(0.4),
         "[Switch to browser for live demo — or embed screen recording here]",
         size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
slide_number(sl, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SLIDE 12 — CONCLUSION + FUTURE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, BG_DARK)
section_header(sl, "Conclusion & Future Directions")

# Key takeaways
take_box = add_rect(sl, Inches(0.5), Inches(1.8), Inches(6), Inches(3.2), BG_CARD, border_color=ACCENT, radius=0.06)
add_text(sl, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.4),
         "What We Built", size=18, color=ACCENT, bold=True)
takeaways = [
    "Real-time fatigue detection from a standard webcam",
    "478-point landmark tracking under 30ms per frame",
    "Adaptive EAR thresholds calibrated per user",
    "DINOv2 + XGBoost fusion for robust classification",
    "Safety-critical alert engine with lock mechanism",
    "Full-stack React 19 dashboard with live visualizations",
]
for j, t in enumerate(takeaways):
    add_text(sl, Inches(0.8), Inches(2.4 + j * 0.4), Inches(5.5), Inches(0.35),
             f"   {t}", size=13, color=LIGHT_GRAY)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(2.52 + j * 0.4), Pt(7), Pt(7))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN
    dot.line.fill.background()

# Future directions
future_box = add_rect(sl, Inches(6.8), Inches(1.8), Inches(6), Inches(3.2), BG_CARD, border_color=PURPLE, radius=0.06)
add_text(sl, Inches(7.1), Inches(1.9), Inches(5.5), Inches(0.4),
         "Future Directions", size=18, color=PURPLE, bold=True)
futures = [
    "LLM-augmented context-aware alert messages",
    "iOS mobile deployment via HTTPS camera",
    "ONNX export + INT8 quantization for edge",
    "Head pose + gaze tracking signals",
    "Longitudinal fatigue trend analytics",
    "Formal validation study (target: 0.90 F1)",
]
for j, t in enumerate(futures):
    add_text(sl, Inches(7.1), Inches(2.4 + j * 0.4), Inches(5.5), Inches(0.35),
             f"   {t}", size=13, color=LIGHT_GRAY)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), Inches(2.52 + j * 0.4), Pt(7), Pt(7))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PURPLE
    dot.line.fill.background()

# GitHub callout
gh_box = add_rect(sl, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.8), RGBColor(0x0A, 0x2E, 0x1A), border_color=GREEN, radius=0.05)
add_text(sl, Inches(0.8), Inches(5.35), Inches(11.5), Inches(0.3),
         "Source Code", size=14, color=GREEN, bold=True)
add_text(sl, Inches(0.8), Inches(5.65), Inches(11.5), Inches(0.3),
         "github.com/omunite215/Project_Pilotguard", size=15, color=ACCENT2)

# Thank you
add_text(sl, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
         "Thank you!", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

slide_number(sl, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "PilotGuard_Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {TOTAL_SLIDES}")
